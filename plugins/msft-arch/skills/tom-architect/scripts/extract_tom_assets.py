#!/usr/bin/env python3
"""
TOM Asset Extraction Pipeline

Extracts structured content from PPTX files (primary) with PDF table verification.
Outputs markdown reference files organized by domain with YAML frontmatter.

Usage:
    pip install -r requirements.txt
    python extract_tom_assets.py --input /path/to/tom-assets --output /path/to/references/domains
"""

import argparse
import os
import re
import sys
from pathlib import Path
from datetime import date

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx>=0.6.21")
    sys.exit(1)

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
    print("WARNING: pdfplumber not installed. PDF table verification disabled.")

# Brand terms to remove/replace (vendor-neutral)
BRAND_TERMS = [
    "Powered Enterprise", "powered enterprise", "Powered enterprise",
    "Connected Enterprise",
]

BRAND_REPLACEMENTS = {
    "Powered Enterprise": "Enterprise Transformation Framework",
    "powered enterprise": "enterprise transformation framework",
    "Connected Enterprise": "Connected Enterprise Framework",
}

# Section metadata
SECTION_INFO = {
    "0.1": {"slug": "intro-powered", "layer": "intro", "title": "Framework Introduction"},
    "0.2": {"slug": "intro-tom", "layer": "intro", "title": "TOM Introduction"},
    "1.1": {"slug": "process-taxonomies", "layer": "processes", "title": "Process Taxonomies"},
    "1.2": {"slug": "maturity-models", "layer": "processes", "title": "Maturity Models"},
    "1.3": {"slug": "process-flows", "layer": "processes", "title": "Role-Based Process Flows"},
    "1.4": {"slug": "leading-practices", "layer": "processes", "title": "Leading Practices & Design Considerations"},
    "2.1": {"slug": "process-owners", "layer": "organization", "title": "Global Process Owners Overlay"},
    "2.2": {"slug": "role-mapping", "layer": "organization", "title": "Position to Role Mapping & Sizing"},
    "2.3": {"slug": "job-profiles", "layer": "organization", "title": "Functional Position Job Profiles"},
    "3.1": {"slug": "service-delivery", "layer": "service", "title": "Service Delivery Model Overlay"},
    "3.2": {"slug": "service-management", "layer": "service", "title": "Service Management Framework"},
    "4.1": {"slug": "ai-augmentation", "layer": "technology", "title": "AI Augmentation Overlay"},
    "4.2": {"slug": "technology-overlay", "layer": "technology", "title": "Supporting Technology Overlay"},
    "4.3": {"slug": "app-architecture", "layer": "technology", "title": "Application Architecture & Data Flows"},
    "4.4": {"slug": "environment-arch", "layer": "technology", "title": "Environment Architecture"},
    "5.1": {"slug": "kpis-benchmarks", "layer": "data-analytics", "title": "KPIs Linked to Benchmarks"},
    "5.2": {"slug": "reporting-dashboards", "layer": "data-analytics", "title": "Reporting Package & Dashboards"},
    "5.3": {"slug": "mdm-governance", "layer": "data-analytics", "title": "MDM Design & Governance"},
    "6.1": {"slug": "security-controls", "layer": "governance", "title": "Security & Controls"},
    "6.2": {"slug": "policies", "layer": "governance", "title": "Policies"},
}

DOMAIN_MAP = {
    "Finance": "finance",
    "HR": "hr",
    "Procurement": "procurement",
    "SCM": "scm",
    "Cyber": "cyber",
    "Sustainability": "sustainability",
}


def remove_brand(text: str) -> str:
    """Remove brand references from text."""
    for original, replacement in BRAND_REPLACEMENTS.items():
        text = text.replace(original, replacement)
    for term in BRAND_TERMS:
        text = text.replace(term, "")
    # Clean up double spaces and empty parentheses
    text = re.sub(r"\(\s*\)", "", text)
    text = re.sub(r"\s{2,}", " ", text)
    text = re.sub(r"^\s*[-–—]\s*$", "", text, flags=re.MULTILINE)
    return text.strip()


def parse_filename(filename: str):
    """Parse TOM filename to extract section and domain."""
    name = Path(filename).stem
    # Match patterns like "1.1 Process Taxonomies - Finance" or "4.1 AI Augmentation Overlay-Cyber"
    match = re.match(r"^(\d+\.\d+)\s+(.+?)[\s]*[-–][\s]*(\w+)$", name)
    if match:
        section_id = match.group(1)
        domain_raw = match.group(3).strip()
        return section_id, domain_raw
    # Try intro files "0.1 Powered Intro" or "0.2 TOM Intro"
    match = re.match(r"^(\d+\.\d+)\s+(.+)$", name)
    if match:
        section_id = match.group(1)
        return section_id, None
    return None, None


def extract_pptx(filepath: str) -> dict:
    """Extract structured content from a PPTX file."""
    prs = Presentation(filepath)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "number": slide_num,
            "title": "",
            "body": [],
            "tables": [],
            "notes": "",
        }

        # Extract title
        if slide.shapes.title:
            slide_info["title"] = remove_brand(slide.shapes.title.text.strip())

        # Extract body text and tables
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != slide_info["title"]:
                        cleaned = remove_brand(text)
                        if cleaned:
                            # Detect bullet level
                            level = para.level if para.level else 0
                            prefix = "  " * level + "- " if level > 0 else ""
                            slide_info["body"].append(f"{prefix}{cleaned}")

            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(remove_brand(cell.text.strip()))
                    table_data.append(row_data)
                if table_data:
                    slide_info["tables"].append(table_data)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["notes"] = remove_brand(notes_text)

        slides_data.append(slide_info)

    return {"slides": slides_data, "slide_count": len(slides_data)}


def extract_pdf_tables(filepath: str) -> list:
    """Extract tables from PDF for verification."""
    if pdfplumber is None:
        return []
    tables = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_tables = page.extract_tables()
                if page_tables:
                    for t in page_tables:
                        cleaned = []
                        for row in t:
                            cleaned.append([remove_brand(cell or "") for cell in row])
                        tables.append(cleaned)
    except Exception as e:
        print(f"  WARNING: PDF table extraction failed for {filepath}: {e}")
    return tables


def table_to_markdown(table: list) -> str:
    """Convert a table (list of rows) to markdown format."""
    if not table or not table[0]:
        return ""
    # Header row
    header = table[0]
    md = "| " + " | ".join(str(h) for h in header) + " |\n"
    md += "| " + " | ".join("---" for _ in header) + " |\n"
    # Data rows
    for row in table[1:]:
        # Pad row to match header length
        padded = list(row) + [""] * (len(header) - len(row))
        md += "| " + " | ".join(str(c) for c in padded[:len(header)]) + " |\n"
    return md


def slides_to_markdown(data: dict, section_id: str, domain: str, section_info: dict) -> str:
    """Convert extracted slide data to structured markdown."""
    title = section_info["title"]
    layer = section_info["layer"]
    slug = section_info["slug"]
    today = date.today().isoformat()

    # YAML frontmatter
    md = f"""---
section: "{section_id}"
domain: "{domain}"
layer: "{layer}"
title: "{title}"
loading_priority: 2
keywords: [{slug}, {domain}, {layer}, tom, target-operating-model]
version: "1.0"
last_updated: "{today}"
---

# {section_id} {title}: {domain.title()}

"""

    for slide in data["slides"]:
        # Skip empty slides
        if not slide["title"] and not slide["body"] and not slide["tables"]:
            continue

        # Slide title as heading
        if slide["title"]:
            md += f"## {slide['title']}\n\n"

        # Body text
        for line in slide["body"]:
            md += f"{line}\n"
        if slide["body"]:
            md += "\n"

        # Tables
        for table in slide["tables"]:
            md += table_to_markdown(table) + "\n"

        # Speaker notes as blockquote
        if slide["notes"]:
            md += f"> **Architect Notes**: {slide['notes']}\n\n"

    return md


def generate_domain_summary(domain: str, extracted_sections: list, output_dir: str):
    """Generate a _domain-summary.md file for a domain."""
    today = date.today().isoformat()
    md = f"""---
domain: "{domain}"
type: "domain-summary"
loading_priority: 1
version: "1.0"
last_updated: "{today}"
---

# {domain.title()} Domain: TOM Reference Index

## Available Sections

| Section | Title | File |
|---------|-------|------|
"""
    for section_id in sorted(extracted_sections):
        info = SECTION_INFO.get(section_id, {})
        slug = info.get("slug", section_id)
        title = info.get("title", "Unknown")
        md += f"| {section_id} | {title} | `{section_id}-{slug}.md` |\n"

    md += f"\n## Domain Coverage\n\n"
    md += f"This domain has **{len(extracted_sections)}** TOM sections extracted.\n"
    md += f"Load individual section files for detailed process taxonomies, maturity models, "
    md += f"role mappings, technology overlays, KPIs, and governance frameworks.\n"

    filepath = os.path.join(output_dir, domain, "_domain-summary.md")
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(md)
    print(f"  Generated: {filepath}")


def main():
    parser = argparse.ArgumentParser(description="Extract TOM assets to structured markdown")
    parser.add_argument("--input", required=True, help="Path to tom-assets folder")
    parser.add_argument("--output", required=True, help="Path to output references/domains folder")
    parser.add_argument("--pdf-verify", action="store_true", help="Enable PDF table verification")
    args = parser.parse_args()

    input_dir = Path(args.input)
    output_dir = Path(args.output)

    if not input_dir.exists():
        print(f"ERROR: Input directory not found: {input_dir}")
        sys.exit(1)

    # Find all PPTX files
    pptx_files = sorted(input_dir.glob("*.pptx"))
    print(f"Found {len(pptx_files)} PPTX files")

    # Track extracted sections per domain
    domain_sections = {d: [] for d in DOMAIN_MAP.values()}
    stats = {"processed": 0, "skipped": 0, "errors": 0}

    for pptx_path in pptx_files:
        section_id, domain_raw = parse_filename(pptx_path.name)

        if section_id is None:
            print(f"  SKIP (unparseable): {pptx_path.name}")
            stats["skipped"] += 1
            continue

        if section_id not in SECTION_INFO:
            print(f"  SKIP (unknown section {section_id}): {pptx_path.name}")
            stats["skipped"] += 1
            continue

        section_info = SECTION_INFO[section_id]

        # Handle intro files (no domain)
        if domain_raw is None:
            domain = "intro"
        else:
            domain = DOMAIN_MAP.get(domain_raw)
            if domain is None:
                print(f"  SKIP (unknown domain '{domain_raw}'): {pptx_path.name}")
                stats["skipped"] += 1
                continue

        print(f"Processing: {pptx_path.name} → {domain}/{section_id}-{section_info['slug']}.md")

        try:
            # Extract PPTX content
            data = extract_pptx(str(pptx_path))

            # Optional PDF verification
            if args.pdf_verify and pdfplumber:
                pdf_name = pptx_path.stem + ".pdf"
                pdf_path = input_dir / pdf_name
                if pdf_path.exists():
                    pdf_tables = extract_pdf_tables(str(pdf_path))
                    if pdf_tables:
                        print(f"  PDF verification: {len(pdf_tables)} tables found")

            # Convert to markdown
            md_content = slides_to_markdown(data, section_id, domain, section_info)

            # Ensure output directory exists
            if domain == "intro":
                out_dir = output_dir / "_intro"
            else:
                out_dir = output_dir / domain
            out_dir.mkdir(parents=True, exist_ok=True)

            # Write output file
            out_file = out_dir / f"{section_id}-{section_info['slug']}.md"
            with open(out_file, "w", encoding="utf-8") as f:
                f.write(md_content)

            # Track for domain summary
            if domain in domain_sections:
                domain_sections[domain].append(section_id)

            stats["processed"] += 1

        except Exception as e:
            print(f"  ERROR: {e}")
            stats["errors"] += 1

    # Generate domain summaries
    print("\nGenerating domain summaries...")
    for domain, sections in domain_sections.items():
        if sections:
            generate_domain_summary(domain, sections, str(output_dir))

    # Print summary
    print(f"\n{'='*60}")
    print(f"Extraction Complete")
    print(f"{'='*60}")
    print(f"Processed: {stats['processed']}")
    print(f"Skipped:   {stats['skipped']}")
    print(f"Errors:    {stats['errors']}")
    print(f"Output:    {output_dir}")


if __name__ == "__main__":
    main()
